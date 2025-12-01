from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_dark_theme_presentation():
    # 1. Create Presentation
    prs = Presentation()
    
    # Set to 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Define Colors
    bg_color = RGBColor(15, 23, 42)      # Dark Slate Background
    title_color = RGBColor(241, 245, 249) # White-ish Title
    text_color = RGBColor(203, 213, 225)  # Light Grey Text
    accent_color = RGBColor(56, 189, 248) # Cyan Accent

    # --- Helper Function to apply Theme ---
    def apply_theme(slide, title_text, content_text=None):
        # Set Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # Set Title
        if slide.shapes.title:
            title = slide.shapes.title
            title.text = title_text
            
            # Title Style
            tf = title.text_frame
            p = tf.paragraphs[0]
            p.font.name = "Urbanist" # Fallback for Urbanist
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = title_color
            p.alignment = PP_ALIGN.LEFT

        # Set Content (if exists)
        if content_text and len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = content_text
            
            # Apply style to all paragraphs
            for p in tf.paragraphs:
                p.font.name = "Urbanist"
                p.font.size = Pt(22)
                p.font.color.rgb = text_color
                p.space_after = Pt(14)

    # --- SLIDE 1: Title Slide ---
    slide_layout = prs.slide_layouts[0] # Title Slide Layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Custom Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    title = slide.shapes.title
    title.text = "Module 1: Processes and Threads"
    title.text_frame.paragraphs[0].font.name = "Urbanist"
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Foundations of System Architecture & Concurrency"
    subtitle.text_frame.paragraphs[0].font.name = "Urbanist"
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_color

    # --- SLIDE 2: Understanding Processes ---
    slide_layout = prs.slide_layouts[1] # Title + Content
    slide = prs.slides.add_slide(slide_layout)
    content = (
        "Core Concepts\n"
        "- Process vs. Application: One application can spawn multiple independent processes.\n"
        "- Lifecycle States: Creation, Running, Waiting, Terminated.\n"
        "- Memory Protection: Every process operates in isolated memory space."
    )
    apply_theme(slide, "Understanding Processes", content)

    # --- SLIDE 3: Why It Matters ---
    slide = prs.slides.add_slide(slide_layout)
    content = (
        "Process Isolation & Safety\n"
        "- Your server application runs as a distinct process.\n"
        "- Isolation allows the OS to contain crashes.\n"
        "- Critical for debugging memory leaks and segmentation faults."
    )
    apply_theme(slide, "Why It Matters", content)

    # --- SLIDE 4: Threads vs Processes ---
    slide = prs.slides.add_slide(slide_layout)
    content = (
        "Key Distinctions\n"
        "- Threads share the same memory space within a process.\n"
        "- Resource Sharing: Threads share code, data, and files.\n"
        "- Overhead: Threads are 'lightweight'; faster creation than processes.\n"
        "- Risk: A bug in one thread can crash the entire process."
    )
    apply_theme(slide, "Threads vs. Processes", content)

    # --- SLIDE 5: Synchronization Challenges ---
    slide = prs.slides.add_slide(slide_layout)
    content = (
        "- Race Conditions: Simultaneous access leads to unpredictable outcomes.\n"
        "- Deadlocks: Threads wait indefinitely for each other to release resources.\n"
        "- Context Switching: CPU overhead when switching threads."
    )
    apply_theme(slide, "Synchronization Challenges", content)

    # --- SLIDE 6: Language Specifics ---
    slide = prs.slides.add_slide(slide_layout)
    content = (
        "Java: JVM threads map to OS threads (ExecutorService).\n"
        "Python: Limited by GIL (Global Interpreter Lock). Use multiprocessing for CPU tasks.\n"
        "Node.js: Single-threaded Event Loop. Use Worker Threads for CPU tasks."
    )
    apply_theme(slide, "Language Specifics", content)

    # --- SLIDE 7: Practical Examples ---
    slide = prs.slides.add_slide(slide_layout)
    content = (
        "- Thread Pool Sizing: Balance pool size to match CPU cores.\n"
        "- Async vs Threading: Async for I/O; Threading for CPU tasks.\n"
        "- Monitoring: Track thread states and deadlock detection."
    )
    apply_theme(slide, "Practical Examples", content)

    # --- SLIDE 8: Q&A ---
    slide_layout = prs.slide_layouts[0] # Title Slide for End
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    title = slide.shapes.title
    title.text = "Questions?"
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Deep dive into process management"
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_color

    # Save
    prs.save('Module1_Processes.pptx')
    print("Presentation saved successfully as Module1_Processes.pptx")

if __name__ == "__main__":
    create_dark_theme_presentation()